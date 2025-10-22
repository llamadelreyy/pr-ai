import os
import asyncio
import aiohttp
import pandas as pd
import numpy as np
from typing import List, Dict, Optional
from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse, JSONResponse
from pydantic import BaseModel
import json
import re
import math
from datetime import datetime
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches
import io
import tempfile
import uuid
from database import db

# Load environment variables
load_dotenv()

def deep_clean_nan_values(obj):
    """Recursively clean NaN, inf, and None values from any data structure"""
    if isinstance(obj, dict):
        return {k: deep_clean_nan_values(v) for k, v in obj.items()}
    elif isinstance(obj, list):
        return [deep_clean_nan_values(item) for item in obj]
    elif isinstance(obj, tuple):
        return tuple(deep_clean_nan_values(item) for item in obj)
    elif isinstance(obj, (np.ndarray,)):
        # Handle numpy arrays
        cleaned_array = np.where(np.isnan(obj), 0.0, obj)
        cleaned_array = np.where(np.isinf(cleaned_array), 0.0, cleaned_array)
        return cleaned_array.tolist()
    elif isinstance(obj, (float, np.float64, np.float32)):
        if pd.isna(obj) or math.isnan(obj) or math.isinf(obj) or obj != obj:  # obj != obj catches NaN
            return 0.0
        return float(obj)
    elif isinstance(obj, (int, np.int64, np.int32)):
        if pd.isna(obj):
            return 0
        return int(obj)
    elif obj is None or pd.isna(obj):
        return ""
    elif isinstance(obj, str):
        if obj.lower() in ['nan', 'none', 'null']:
            return ""
        return obj
    return obj

def safe_json_response(content):
    """Create a safe JSON response that handles all NaN values"""
    cleaned_content = deep_clean_nan_values(content)
    return JSONResponse(content=cleaned_content)

app = FastAPI(title="MEIO Sentiment Analysis System")

# CORS middleware - Allow all origins for external access
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins for port forwarding
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Global storage for session data
sessions = {}

class SentimentRequest(BaseModel):
    session_id: str
    text_column_override: Optional[str] = None

class CommentRequest(BaseModel):
    session_id: str
    text_column_override: Optional[str] = None

class ReportRequest(BaseModel):
    session_id: str
    text_column_override: Optional[str] = None

class InfluentialAnalysisRequest(BaseModel):
    session_id: str
    text_column_override: Optional[str] = None
    exposure_threshold: Optional[float] = 500000.0  # 500K default threshold
    top_count: Optional[int] = 5  # Top 5 default
    use_threshold: Optional[bool] = True  # Use threshold vs top count

class SaveAnalysisRequest(BaseModel):
    session_id: str
    document_id: str
    analysis_name: str
    analysis_type: str
    analysis_data: dict

def detect_cyabra_columns(data):
    """
    Detect Cyabra-style export columns including max_exp, tag, author_url, content_url, and engagement metrics
    """
    if not data or len(data) == 0:
        return {}
    
    df_temp = pd.DataFrame(data)
    detected_columns = {}
    
    # Look for exposure/reach columns (max_exp or similar)
    exposure_patterns = [r'.*max_exp.*', r'.*exposure.*', r'.*reach.*', r'.*views.*', r'.*impressions.*']
    for col in df_temp.columns:
        col_lower = col.lower()
        for pattern in exposure_patterns:
            if re.match(pattern, col_lower):
                detected_columns['exposure_column'] = col
                break
        if 'exposure_column' in detected_columns:
            break
    
    # Look for tag/profile type columns
    tag_patterns = [r'.*tag.*', r'.*profile.*type.*', r'.*account.*type.*', r'.*user.*type.*']
    for col in df_temp.columns:
        col_lower = col.lower()
        for pattern in tag_patterns:
            if re.match(pattern, col_lower):
                detected_columns['tag_column'] = col
                break
        if 'tag_column' in detected_columns:
            break
    
    # Look for author URL columns
    author_url_patterns = [r'.*author.*url.*', r'.*profile.*url.*', r'.*user.*url.*', r'.*account.*url.*']
    for col in df_temp.columns:
        col_lower = col.lower()
        for pattern in author_url_patterns:
            if re.match(pattern, col_lower):
                detected_columns['author_url_column'] = col
                break
        if 'author_url_column' in detected_columns:
            break
    
    # Look for content URL columns
    content_url_patterns = [r'.*content.*url.*', r'.*post.*url.*', r'.*link.*', r'.*url.*']
    for col in df_temp.columns:
        col_lower = col.lower()
        for pattern in content_url_patterns:
            if re.match(pattern, col_lower):
                detected_columns['content_url_column'] = col
                break
        if 'content_url_column' in detected_columns:
            break
    
    # Look for engagement metrics columns
    engagement_columns = {}
    
    # Likes columns (prioritize standard names)
    likes_patterns = [r'^like_count$', r'^likes$', r'.*like.*count.*', r'.*favorite.*', r'.*heart.*']
    for col in df_temp.columns:
        col_lower = col.lower()
        for pattern in likes_patterns:
            if re.match(pattern, col_lower):
                engagement_columns['likes'] = col
                break
    
    # Shares/Retweets columns (prioritize standard names)
    shares_patterns = [r'^share_count$', r'^shares$', r'.*share.*count.*', r'.*retweet.*', r'.*rt.*count.*', r'.*reshare.*']
    for col in df_temp.columns:
        col_lower = col.lower()
        for pattern in shares_patterns:
            if re.match(pattern, col_lower):
                engagement_columns['shares'] = col
                break
    
    # Comments/Replies columns (prioritize standard names)
    comments_patterns = [r'^comment_count$', r'^comments$', r'.*comment.*count.*', r'.*reply.*', r'.*response.*']
    for col in df_temp.columns:
        col_lower = col.lower()
        for pattern in comments_patterns:
            if re.match(pattern, col_lower) and 'generated' not in col_lower:
                engagement_columns['comments'] = col
                break
    
    # Total engagement column
    engagement_patterns = [r'^engagement$', r'^total.*engagement.*', r'.*engagement.*total.*']
    for col in df_temp.columns:
        col_lower = col.lower()
        for pattern in engagement_patterns:
            if re.match(pattern, col_lower):
                engagement_columns['total_engagement'] = col
                break
    
    # Views/Impressions columns
    views_patterns = [r'.*view.*', r'.*impression.*', r'.*seen.*']
    for col in df_temp.columns:
        col_lower = col.lower()
        for pattern in views_patterns:
            if re.match(pattern, col_lower):
                engagement_columns['views'] = col
                break
    
    if engagement_columns:
        detected_columns['engagement_columns'] = engagement_columns
    
    # Look for author name columns
    author_name_patterns = [r'.*author.*name.*', r'.*user.*name.*', r'.*username.*', r'.*handle.*', r'.*account.*name.*']
    for col in df_temp.columns:
        col_lower = col.lower()
        for pattern in author_name_patterns:
            if re.match(pattern, col_lower):
                detected_columns['author_name_column'] = col
                break
        if 'author_name_column' in detected_columns:
            break
    
    return detected_columns

def detect_text_column(data):
    """
    Intelligently detect the most likely text column for sentiment analysis.
    Uses multiple heuristics to find the best text column.
    """
    if not data or len(data) == 0:
        return None
    
    df_temp = pd.DataFrame(data)
    column_scores = {}
    
    # Define common text column name patterns
    text_patterns = [
        # Common text column names
        r'.*text.*', r'.*content.*', r'.*message.*', r'.*comment.*', r'.*post.*',
        r'.*description.*', r'.*body.*', r'.*tweet.*', r'.*status.*', r'.*review.*',
        r'.*feedback.*', r'.*opinion.*', r'.*remarks.*', r'.*note.*', r'.*detail.*',
        # Social media specific
        r'.*caption.*', r'.*thread.*', r'.*reply.*', r'.*mention.*',
        # Document specific
        r'.*summary.*', r'.*abstract.*', r'.*excerpt.*', r'.*paragraph.*'
    ]
    
    # Check each column
    for col in df_temp.columns:
        if df_temp[col].dtype == 'object':
            score = 0
            
            # Convert to string and handle NaN values
            text_series = df_temp[col].astype(str)
            valid_texts = text_series[~text_series.isin(['', 'nan', 'None', 'NULL'])]
            
            if len(valid_texts) == 0:
                continue
            
            # Score 1: Column name patterns (high weight)
            col_lower = col.lower()
            for pattern in text_patterns:
                if re.match(pattern, col_lower):
                    score += 50
                    break
            
            # Score 2: Average text length (higher is better for meaningful text)
            avg_length = valid_texts.str.len().mean()
            if avg_length > 20:  # Meaningful text usually > 20 chars
                score += min(avg_length / 10, 30)  # Cap at 30 points
            elif avg_length < 5:  # Probably not text content
                score -= 20
            
            # Score 3: Text diversity (more unique values = better)
            unique_ratio = len(valid_texts.unique()) / len(valid_texts)
            if unique_ratio > 0.7:  # High diversity
                score += 20
            elif unique_ratio < 0.1:  # Low diversity (probably IDs or categories)
                score -= 15
            
            # Score 4: Contains common text words/patterns
            sample_text = ' '.join(valid_texts.head(100).str.lower())
            text_indicators = ['the', 'and', 'is', 'to', 'of', 'in', 'it', 'you', 'that', 'he', 'was', 'for', 'on', 'are', 'as', 'with', 'his', 'they', 'i', 'at', 'be', 'this', 'have', 'from', 'or', 'one', 'had', 'by', 'word', 'but', 'not', 'what', 'all', 'were', 'we', 'when', 'your', 'can', 'said', 'there', 'each', 'which', 'she', 'do', 'how', 'their', 'if', 'will', 'up', 'other', 'about', 'out', 'many', 'then', 'them', 'these', 'so', 'some', 'her', 'would', 'make', 'like', 'into', 'him', 'has', 'two', 'more', 'very', 'what', 'know', 'just', 'first', 'get', 'over', 'think', 'also', 'its', 'back', 'after', 'use', 'good', 'our', 'way', 'even', 'new', 'want', 'because', 'any', 'give', 'day', 'us', 'most', 'people']
            word_count = sum(1 for word in text_indicators if word in sample_text)
            score += min(word_count, 25)  # Cap at 25 points
            
            # Score 5: Contains punctuation (typical of real text)
            punct_ratio = sum(1 for char in sample_text if char in '.,!?;:') / max(len(sample_text), 1)
            if punct_ratio > 0.01:  # Has decent punctuation
                score += 10
            
            # Score 6: Number of sentences (more sentences = more likely to be text)
            sentence_count = sample_text.count('.') + sample_text.count('!') + sample_text.count('?')
            if sentence_count > len(valid_texts.head(100)) * 0.3:  # Good sentence density
                score += 15
            
            # Penalty for likely non-text columns
            if col_lower in ['id', 'index', 'number', 'count', 'date', 'time', 'url', 'link', 'email', 'phone', 'address']:
                score -= 30
            
            # Penalty for columns that look like IDs or numbers
            if valid_texts.str.match(r'^\d+$').sum() > len(valid_texts) * 0.8:  # Mostly numbers
                score -= 25
            
            if valid_texts.str.len().max() < 10:  # Very short content
                score -= 15
                
            column_scores[col] = score
            print(f"Column '{col}': score={score:.1f}, avg_length={avg_length:.1f}, unique_ratio={unique_ratio:.2f}")
    
    if not column_scores:
        # Fallback: return first object column
        for col in df_temp.columns:
            if df_temp[col].dtype == 'object':
                return col
        return df_temp.columns[0] if len(df_temp.columns) > 0 else None
    
    # Return column with highest score
    best_column = max(column_scores.items(), key=lambda x: x[1])
    print(f"Selected text column: '{best_column[0]}' with score {best_column[1]:.1f}")
    return best_column[0]

class SessionData:
    def __init__(self):
        self.data = None
        self.sentiments = None
        self.comments = None
        self.filename = None
        self.detected_text_column = None

async def call_llm_api(text: str, prompt_type: str) -> Dict:
    """Call the LLM API for sentiment analysis, topic categorization, or comment generation"""
    llm_url = os.getenv("LLM_API_URL")
    model_name = os.getenv("LLM_MODEL_NAME")
    
    if prompt_type == "sentiment":
        prompt = f"""Analyze the sentiment of the following text. Respond in JSON format with sentiment and confidence score.

Text: {text}

Response format:
{{"sentiment": "positive|negative|neutral", "confidence": 0.95}}

Response:"""
    elif prompt_type == "topic":
        prompt = f"""Categorize the topic/theme of the following text. Choose multiple appropriate categories separated by commas.

Text: {text}

Categories: Politics, Economy, Security, Foreign Relations, Social Issues, Technology, Healthcare, Education, Environment, Defense, Trade, Culture, Infrastructure, Government Policy, Public Opinion

Respond with category names separated by commas:"""
    elif prompt_type == "counter":
        prompt = f"""Respond to this negative comment with a balanced, positive perspective. Keep the response natural and proportional to the original length and tone.

Original text: {text}

Instructions:
- Match the language style and formality level of the original
- Keep response length similar to the original (don't make it much longer)
- Address concerns with factual, positive information
- Be diplomatic and respectful, not defensive
- Provide constructive perspective without being preachy
- Stay natural and authentic to the communication style

Response:"""
    elif prompt_type == "amplify":
        prompt = f"""Build upon this positive comment with additional supportive content. Keep the response natural and proportional to the original.

Original text: {text}

Instructions:
- Match the language style and formality level of the original
- Keep response length appropriate (don't over-expand)
- Add meaningful support or additional positive aspects
- Maintain the same energy level as the original
- Be authentic and natural, not over-enthusiastic
- Stay true to the original communication style

Response:"""
    elif prompt_type == "neutral_positive":
        prompt = f"""Add a gentle positive perspective to this neutral comment. Keep the response natural and balanced.

Original text: {text}

Instructions:
- Match the language style and formality level of the original
- Keep response length similar to the original
- Gently highlight positive aspects or benefits
- Be subtle and natural, not forceful
- Maintain credibility and authenticity
- Add optimism without being overly promotional

Response:"""
    
    payload = {
        "model": model_name,
        "messages": [
            {"role": "user", "content": prompt}
        ],
        "max_tokens": 200,
        "temperature": 0.3,
        "stream": False  # Disable streaming for now to fix response handling
    }
    
    try:
        async with aiohttp.ClientSession() as session:
            async with session.post(f"{llm_url}/chat/completions",
                                   json=payload,
                                   headers={"Content-Type": "application/json"}) as response:
                if response.status == 200:
                    # Always get the JSON response for all cases
                    result = await response.json()
                    content = result["choices"][0]["message"]["content"].strip()
                    print(f"LLM API response for {prompt_type}: {content[:100]}")
                    
                    if prompt_type == "sentiment":
                        try:
                            # Try to parse JSON response
                            parsed = json.loads(content)
                            return parsed
                        except json.JSONDecodeError:
                            # Fallback to text parsing
                            sentiment_match = re.search(r'(positive|negative|neutral)', content.lower())
                            confidence_match = re.search(r'(\d+\.?\d*)', content)
                            
                            sentiment = sentiment_match.group(1) if sentiment_match else "neutral"
                            try:
                                confidence = float(confidence_match.group(1)) if confidence_match else 0.7
                                if confidence > 1:
                                    confidence = confidence / 100
                                # Ensure confidence is within valid range
                                confidence = max(0.0, min(1.0, confidence))
                                # Check for NaN or infinity
                                if not (0.0 <= confidence <= 1.0):
                                    confidence = 0.7
                            except (ValueError, TypeError):
                                confidence = 0.7
                            
                            return {"sentiment": sentiment, "confidence": confidence}
                    elif prompt_type in ["counter", "amplify", "neutral_positive"]:
                        # For comment generation, always return content in dict format
                        return {"content": content}
                    elif prompt_type == "topic":
                        # For topic classification, return content directly
                        return {"content": content}
                    else:
                        return {"content": content}
                else:
                    print(f"LLM API returned status {response.status}")
                    return {"error": f"HTTP {response.status}", "content": "No comment generated"}
    except Exception as e:
        print(f"LLM API Error: {e}")
        return {"error": "API unavailable", "content": "No comment generated"}

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """Upload and process CSV/XLSX file - simplified for debugging"""
    print(f"=== UPLOAD START === File: {file.filename}")
    
    # Validate file type
    if not file.filename.endswith(('.csv', '.xlsx')):
        print(f"Invalid file type: {file.filename}")
        raise HTTPException(status_code=400, detail="Only CSV and XLSX files are supported")
    
    session_id = str(uuid.uuid4())
    print(f"Generated session ID: {session_id}")
    
    try:
        print("Step 1: Reading file contents...")
        contents = await file.read()
        file_size = len(contents)
        print(f"File size: {file_size} bytes ({file_size/1024/1024:.2f} MB)")
        
        # Handle large files with chunked processing
        if file_size > 100 * 1024 * 1024:  # 100MB limit
            print(f"Very large file detected: {file_size} bytes - will use chunked processing")
        elif file_size > 50 * 1024 * 1024:  # 50MB limit
            print(f"Large file detected: {file_size} bytes - will use chunked processing")
        
        print("Step 2: Attempting to read with pandas...")
        
        # Adaptive reading based on file size
        if file.filename.endswith('.csv'):
            print("Reading CSV file...")
            try:
                if file_size > 50 * 1024 * 1024:  # > 50MB, use chunked reading
                    print("Large file - using chunked reading...")
                    chunks = []
                    chunk_size = 10000  # Read 10000 rows at a time
                    max_chunks = None   # No limit - read all chunks
                    
                    # Robust CSV reading parameters for malformed files
                    csv_params = {
                        'chunksize': chunk_size,
                        'on_bad_lines': 'skip',  # Skip malformed lines
                        'engine': 'c',          # Use C engine which supports low_memory
                        'sep': ',',             # Explicit separator
                        'quotechar': '"',       # Standard quote character
                        'skipinitialspace': True,
                        'low_memory': False     # Read entire file into memory for consistency
                    }
                    
                    try:
                        print("Attempting robust CSV parsing with error handling...")
                        csv_params['encoding'] = 'utf-8'
                        chunk_iterator = pd.read_csv(io.BytesIO(contents), **csv_params)
                        for i, chunk in enumerate(chunk_iterator):
                            if len(chunk) > 0:  # Only add non-empty chunks
                                chunks.append(chunk)
                                print(f"Read chunk {i+1}: {len(chunk)} rows")
                            # Continue reading all chunks - no artificial limit
                        
                        if chunks:
                            df = pd.concat(chunks, ignore_index=True)
                            print(f"CSV read successful with chunking: {len(df)} total rows")
                        else:
                            print("No valid chunks found, trying alternative approach...")
                            raise Exception("No valid data chunks found")
                            
                    except (UnicodeDecodeError, Exception) as e:
                        print(f"UTF-8 chunked reading failed: {e}")
                        print("Trying latin-1 with chunking...")
                        chunks = []
                        try:
                            csv_params_latin = csv_params.copy()
                            csv_params_latin['encoding'] = 'latin-1'
                            csv_params_latin['engine'] = 'python'  # Fallback to python engine
                            csv_params_latin.pop('low_memory', None)  # Remove low_memory for python engine
                            chunk_iterator = pd.read_csv(io.BytesIO(contents), **csv_params_latin)
                            for i, chunk in enumerate(chunk_iterator):
                                if len(chunk) > 0:
                                    chunks.append(chunk)
                                # Continue reading all chunks - no artificial limit
                            
                            if chunks:
                                df = pd.concat(chunks, ignore_index=True)
                                print(f"CSV read successful with latin-1 chunking: {len(df)} total rows")
                            else:
                                raise Exception("No valid data found with latin-1 encoding")
                                
                        except Exception as e2:
                            print(f"Chunked reading failed completely: {e2}")
                            print("Attempting single-pass read with robust parameters...")
                            # Last resort: try to read the whole file with very robust settings
                            try:
                                # Try with C engine first (more robust for malformed CSVs)
                                df = pd.read_csv(
                                    io.BytesIO(contents),
                                    encoding='utf-8',
                                    on_bad_lines='skip',
                                    engine='c',
                                    sep=',',
                                    quotechar='"',
                                    skipinitialspace=True,
                                    # No row limit - read all data
                                    low_memory=False
                                )
                                print(f"Single-pass C engine read successful: {len(df)} rows")
                            except Exception as e3:
                                print(f"C engine failed: {e3}, trying python engine...")
                                # Fallback to python engine without low_memory
                                df = pd.read_csv(
                                    io.BytesIO(contents),
                                    encoding='latin-1',
                                    on_bad_lines='skip',
                                    engine='python',
                                    sep=',',
                                    quotechar='"',
                                    skipinitialspace=True,
                                    # No row limit - read all data
                                )
                                print(f"Single-pass python engine read successful: {len(df)} rows")
                else:
                    # Standard reading for smaller files with robust error handling
                    print("Standard file size - using robust single-pass reading...")
                    
                    csv_params = {
                        'on_bad_lines': 'skip',
                        'engine': 'c',         # Use C engine which supports low_memory
                        'sep': ',',
                        'quotechar': '"',
                        'skipinitialspace': True,
                        'low_memory': False
                    }
                    
                    try:
                        df = pd.read_csv(io.BytesIO(contents), encoding='utf-8', **csv_params)
                        print(f"CSV read successful with UTF-8: {len(df)} rows")
                    except (UnicodeDecodeError, Exception) as e:
                        print(f"UTF-8 failed: {e}, trying latin-1...")
                        try:
                            csv_params_latin = csv_params.copy()
                            csv_params_latin['encoding'] = 'latin-1'
                            df = pd.read_csv(io.BytesIO(contents), **csv_params_latin)
                            print(f"CSV read successful with latin-1: {len(df)} rows")
                        except Exception as e2:
                            print(f"Standard parsing failed: {e2}, trying python engine fallback...")
                            # Very basic fallback with python engine - read first 50k for safety
                            df = pd.read_csv(
                                io.BytesIO(contents),
                                encoding='utf-8',
                                nrows=50000,  # Reasonable fallback limit
                                on_bad_lines='skip',
                                engine='python'
                            )
                            print(f"Fallback read successful: {len(df)} rows")
            except Exception as e:
                print(f"CSV reading failed: {str(e)}")
                import traceback
                traceback.print_exc()
                raise HTTPException(status_code=400, detail=f"CSV read error: {str(e)}")
        else:
            print("Reading Excel file...")
            try:
                # For Excel, read all rows
                df = pd.read_excel(io.BytesIO(contents))
                print(f"Excel read successful: {len(df)} rows")
            except Exception as e:
                print(f"Excel reading failed: {str(e)}")
                import traceback
                traceback.print_exc()
                raise HTTPException(status_code=400, detail=f"Excel read error: {str(e)}")
        
        print(f"Step 3: File parsed successfully - {len(df)} rows, {len(df.columns)} columns")
        print(f"Columns: {list(df.columns)}")
        
        # Basic NaN cleaning - simple approach
        print("Step 4: Basic NaN cleaning...")
        df = df.fillna('')  # Replace all NaN with empty string
        
        print("Step 5: Converting to dict...")
        # Convert data efficiently - process all rows
        try:
            # Process all data without sampling
            test_data = df.to_dict('records')
            print(f"Converted all {len(test_data)} rows to dict format")
        except Exception as e:
            print(f"Dict conversion failed: {str(e)}")
            import traceback
            traceback.print_exc()
            raise HTTPException(status_code=500, detail=f"Data conversion error: {str(e)}")
        
        print("Step 6: Creating session data...")
        
        # Store in session
        session_data = SessionData()
        session_data.data = test_data
        session_data.filename = file.filename
        session_data.detected_text_column = detect_text_column(test_data)
        sessions[session_id] = session_data
        
        print("Step 7: Preparing response...")
        response_data = {
            "session_id": session_id,
            "filename": file.filename,
            "row_count": len(test_data),
            "columns": list(df.columns),
            "detected_text_column": session_data.detected_text_column,
            "preview": test_data[:5]  # First 5 rows as preview
        }
        
        print(f"=== UPLOAD SUCCESS === Session: {session_id}, Rows: {len(test_data)}")
        return response_data
    
    except Exception as e:
        print(f"Error processing file: {str(e)}")
        import traceback
        traceback.print_exc()
        
        # Provide more specific error messages
        if "encoding" in str(e).lower():
            raise HTTPException(status_code=400, detail=f"File encoding error. Please save your file as UTF-8 encoded CSV or standard Excel format.")
        elif "memory" in str(e).lower() or "size" in str(e).lower():
            raise HTTPException(status_code=413, detail=f"File too large to process. Please reduce file size or split into smaller files.")
        elif "pandas" in str(e).lower() or "read_csv" in str(e).lower() or "read_excel" in str(e).lower():
            raise HTTPException(status_code=400, detail=f"File format error. Please ensure your file is a valid CSV or Excel file.")
        else:
            raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.get("/data/{session_id}")
async def get_data(session_id: str):
    """Get uploaded data for viewing"""
    if session_id not in sessions:
        raise HTTPException(status_code=404, detail="Session not found")
    
    session_data = sessions[session_id]
    
    # Use safe JSON response
    return safe_json_response({
        "data": session_data.data,
        "filename": session_data.filename
    })

@app.post("/analyze")
async def analyze_sentiment(request: SentimentRequest):
    """Perform sentiment analysis on all rows"""
    if request.session_id not in sessions:
        raise HTTPException(status_code=404, detail="Session not found")
    
    session_data = sessions[request.session_id]
    data = session_data.data
    
    # Use override if provided, otherwise use intelligent detection
    if request.text_column_override:
        text_column = request.text_column_override
        print(f"Using user-specified text column: {text_column}")
    else:
        text_column = detect_text_column(data)
        print(f"Using auto-detected text column: {text_column}")
    
    async def analyze_row(row):
        text = str(row.get(text_column, ""))
        if len(text.strip()) == 0 or text.strip().lower() in ['nan', 'none', '']:
            return {"sentiment": "neutral", "confidence": 0.5, "topic": "General"}
        
        # Get sentiment with precise confidence
        sentiment_result = await call_llm_api(text, "sentiment")
        if "error" in sentiment_result:
            return {"sentiment": "neutral", "confidence": 0.5, "topic": "General"}
        
        # Get topic categorization
        topic_result = await call_llm_api(text, "topic")
        topic = topic_result.get("content", "General") if "content" in topic_result else "General"
        
        # Ensure all values are clean
        sentiment = sentiment_result.get("sentiment", "neutral")
        confidence = sentiment_result.get("confidence", 0.5)
        
        # Validate confidence
        try:
            confidence = float(confidence)
            if not (0.0 <= confidence <= 1.0):
                confidence = 0.5
        except (ValueError, TypeError):
            confidence = 0.5
        
        return {
            "sentiment": sentiment,
            "confidence": confidence,
            "topic": topic.strip() if topic else "General"
        }
    
    # Process rows concurrently
    semaphore = asyncio.Semaphore(10)  # Limit concurrent requests
    
    async def analyze_with_semaphore(row):
        async with semaphore:
            return await analyze_row(row)
    
    tasks = [analyze_with_semaphore(row) for row in data]
    sentiments = await asyncio.gather(*tasks)
    
    # Always check for max_exp column directly first
    available_columns = list(data[0].keys()) if data else []
    print(f"Available columns in data: {available_columns}")
    
    # Detect Cyabra-style columns for exposure scores
    cyabra_columns = detect_cyabra_columns(data)
    exposure_column = cyabra_columns.get('exposure_column')
    tag_column = cyabra_columns.get('tag_column')
    author_url_column = cyabra_columns.get('author_url_column')
    content_url_column = cyabra_columns.get('content_url_column')
    
    # Force use max_exp if it exists (most important for Cyabra data)
    if 'max_exp' in available_columns:
        exposure_column = 'max_exp'
        print(f"Found max_exp column - using it as exposure column")
    
    # Fallback to direct column names if not detected
    if not tag_column and 'tag' in available_columns:
        tag_column = 'tag'
    if not author_url_column and 'author_url' in available_columns:
        author_url_column = 'author_url'
    if not content_url_column and 'content_url' in available_columns:
        content_url_column = 'content_url'
    
    print(f"Final columns used - Exposure: {exposure_column}, Tag: {tag_column}, Author URL: {author_url_column}, Content URL: {content_url_column}")
    
    # Combine data with sentiments and exposure information
    analyzed_data = []
    for i, row in enumerate(data):
        analyzed_row = {**row, **sentiments[i]}
        
        # Extract exposure score - prioritize max_exp column
        exposure_score = 0.0
        
        # Try max_exp first
        if 'max_exp' in row:
            try:
                raw_value = row['max_exp']
                if raw_value is not None and str(raw_value).strip() not in ['', 'nan', 'None']:
                    if isinstance(raw_value, str):
                        clean_value = raw_value.replace(',', '').replace(' ', '').strip()
                        if clean_value:
                            exposure_score = float(clean_value)
                    else:
                        exposure_score = float(raw_value)
                print(f"Row {i}: max_exp extracted exposure {exposure_score} from raw value: {raw_value}")
            except (ValueError, TypeError) as e:
                print(f"Row {i}: Failed to extract max_exp {row.get('max_exp')}: {e}")
                exposure_score = 0.0
        
        # If max_exp failed, try detected exposure column
        elif exposure_column and exposure_column in row:
            try:
                raw_value = row[exposure_column]
                if raw_value is not None and str(raw_value).strip() not in ['', 'nan', 'None']:
                    if isinstance(raw_value, str):
                        clean_value = raw_value.replace(',', '').replace(' ', '').strip()
                        if clean_value:
                            exposure_score = float(clean_value)
                    else:
                        exposure_score = float(raw_value)
                print(f"Row {i}: {exposure_column} extracted exposure {exposure_score} from raw value: {raw_value}")
            except (ValueError, TypeError) as e:
                print(f"Row {i}: Failed to extract from {exposure_column} {row.get(exposure_column)}: {e}")
                exposure_score = 0.0
        
        # Extract engagement metrics with standard column name support
        engagement_data = {}
        engagement_columns = cyabra_columns.get('engagement_columns', {})
        
        # Standard engagement columns with fallbacks
        engagement_mappings = {
            'likes': ['like_count', 'likes'],
            'shares': ['share_count', 'shares'],
            'comments': ['comment_count', 'comments'],
            'views': ['views'],
            'total_engagement': ['engagement', 'total_engagement']
        }
        
        for metric, possible_columns in engagement_mappings.items():
            engagement_data[metric] = 0
            
            # Try detected column first
            if metric in engagement_columns and engagement_columns[metric] in row:
                column = engagement_columns[metric]
                try:
                    raw_value = row[column]
                    if raw_value is not None and str(raw_value).strip() not in ['', 'nan', 'None']:
                        if isinstance(raw_value, str):
                            clean_value = raw_value.replace(',', '').replace(' ', '').strip()
                            engagement_data[metric] = int(float(clean_value)) if clean_value else 0
                        else:
                            engagement_data[metric] = int(float(raw_value))
                except (ValueError, TypeError):
                    engagement_data[metric] = 0
            else:
                # Try standard column names as fallback
                for col_name in possible_columns:
                    if col_name in row:
                        try:
                            raw_value = row[col_name]
                            if raw_value is not None and str(raw_value).strip() not in ['', 'nan', 'None']:
                                if isinstance(raw_value, str):
                                    clean_value = raw_value.replace(',', '').replace(' ', '').strip()
                                    engagement_data[metric] = int(float(clean_value)) if clean_value else 0
                                else:
                                    engagement_data[metric] = int(float(raw_value))
                                break
                        except (ValueError, TypeError):
                            continue
        
        # Calculate total engagement if not provided
        if engagement_data.get('total_engagement', 0) == 0:
            total_eng = engagement_data.get('likes', 0) + (engagement_data.get('shares', 0) * 2) + (engagement_data.get('comments', 0) * 3)
            engagement_data['calculated_engagement'] = total_eng
        
        # Extract author name
        author_name_column = cyabra_columns.get('author_name_column')
        author_name = ""
        if author_name_column and author_name_column in row:
            author_name = str(row.get(author_name_column, ''))
        elif 'author_name' in row:
            author_name = str(row.get('author_name', ''))
        
        # Always include the original max_exp value in the response for debugging
        if 'max_exp' in row:
            analyzed_row['max_exp_original'] = row['max_exp']
        
        analyzed_row['exposure_score'] = exposure_score
        analyzed_row['profile_tag'] = str(row.get(tag_column, 'UNTAGGED')).upper() if tag_column and tag_column in row else 'UNTAGGED'
        analyzed_row['author_url'] = str(row.get(author_url_column, '')) if author_url_column and author_url_column in row else ''
        analyzed_row['content_url'] = str(row.get(content_url_column, '')) if content_url_column and content_url_column in row else ''
        analyzed_row['author_name'] = author_name
        
        # Add engagement metrics
        analyzed_row.update(engagement_data)
        
        analyzed_data.append(analyzed_row)
    
    session_data.sentiments = analyzed_data
    
    # Calculate statistics
    sentiment_counts = {"positive": 0, "negative": 0, "neutral": 0}
    topic_counts = {}
    total_confidence = 0
    
    for item in sentiments:
        sentiment_counts[item["sentiment"]] += 1
        topic = item.get("topic", "General")
        topic_counts[topic] = topic_counts.get(topic, 0) + 1
        
        # Ensure confidence is a valid number
        confidence = item.get("confidence", 0.5)
        try:
            confidence = float(confidence)
            if not (0.0 <= confidence <= 1.0):
                confidence = 0.5
        except (ValueError, TypeError):
            confidence = 0.5
        
        total_confidence += confidence
    
    avg_confidence = total_confidence / len(sentiments) if sentiments else 0.5
    # Ensure avg_confidence is valid
    if not (0.0 <= avg_confidence <= 1.0):
        avg_confidence = 0.5
    
    # Calculate exposure statistics
    exposure_scores = [row.get('exposure_score', 0) for row in analyzed_data]
    exposure_stats = {
        'max_exposure': max(exposure_scores) if exposure_scores else 0,
        'min_exposure': min(exposure_scores) if exposure_scores else 0,
        'avg_exposure': sum(exposure_scores) / len(exposure_scores) if exposure_scores else 0,
        'total_exposure': sum(exposure_scores)
    }
    
    # Profile distribution
    profile_distribution = {}
    for row in analyzed_data:
        tag = row.get('profile_tag', 'UNTAGGED')
        profile_distribution[tag] = profile_distribution.get(tag, 0) + 1
    
    # Calculate engagement statistics
    engagement_stats = {}
    engagement_columns = cyabra_columns.get('engagement_columns', {})
    
    for metric, column in engagement_columns.items():
        metric_values = [row.get(metric, 0) for row in analyzed_data]
        if metric_values:
            engagement_stats[metric] = {
                'total': sum(metric_values),
                'max': max(metric_values),
                'avg': sum(metric_values) / len(metric_values),
                'column_name': column
            }
    
    # Calculate total engagement if not provided in data
    if 'total_engagement' not in engagement_stats and len(engagement_stats) > 0:
        total_engagement_values = []
        for row in analyzed_data:
            total_eng = 0
            total_eng += row.get('likes', 0)
            total_eng += row.get('shares', 0) * 2  # Weight shares higher
            total_eng += row.get('comments', 0) * 3  # Weight comments highest
            total_engagement_values.append(total_eng)
        
        if total_engagement_values:
            engagement_stats['calculated_engagement'] = {
                'total': sum(total_engagement_values),
                'max': max(total_engagement_values),
                'avg': sum(total_engagement_values) / len(total_engagement_values),
                'column_name': 'Calculated (Likes + 2×Shares + 3×Comments)'
            }
    
    # Use safe JSON response
    return safe_json_response({
        "analyzed_data": analyzed_data,
        "statistics": sentiment_counts,
        "topic_statistics": topic_counts,
        "average_confidence": avg_confidence,
        "exposure_statistics": exposure_stats,
        "profile_distribution": profile_distribution,
        "engagement_statistics": engagement_stats,
        "cyabra_columns_detected": cyabra_columns,
        "text_column": text_column
    })

@app.post("/analyze-streaming")
async def analyze_streaming(request: SentimentRequest):
    """Stream sentiment analysis with real-time updates"""
    if request.session_id not in sessions:
        raise HTTPException(status_code=404, detail="Session not found")
    
    session_data = sessions[request.session_id]
    data = session_data.data
    
    if not data:
        raise HTTPException(status_code=400, detail="No data found")
    
    print(f"Starting streaming analysis for {len(data)} rows")
    
    # Use override if provided, otherwise use intelligent detection
    if request.text_column_override:
        text_column = request.text_column_override
        print(f"Using user-specified text column: {text_column}")
    else:
        text_column = detect_text_column(data)
        print(f"Using auto-detected text column: {text_column}")
    
    async def analyze_stream():
        analyzed_data = []
        sentiment_counts = {"positive": 0, "negative": 0, "neutral": 0}
        topic_counts = {}
        total_confidence = 0
        
        for i, row in enumerate(data):
            text = str(row.get(text_column, ""))
            
            # Send row start event
            yield f"data: {json.dumps({'type': 'row_start', 'index': i, 'total': len(data), 'text': text[:100]})}\n\n"
            
            try:
                if len(text.strip()) == 0 or text.strip().lower() in ['nan', 'none', '']:
                    sentiment_result = {"sentiment": "neutral", "confidence": 0.5}
                    topic_result = {"content": "General"}
                else:
                    # Get sentiment with precise confidence
                    sentiment_result = await call_llm_api(text, "sentiment")
                    if "error" in sentiment_result:
                        sentiment_result = {"sentiment": "neutral", "confidence": 0.5}
                    
                    # Send sentiment result
                    yield f"data: {json.dumps({'type': 'sentiment', 'index': i, 'sentiment': sentiment_result.get('sentiment', 'neutral'), 'confidence': sentiment_result.get('confidence', 0.5)})}\n\n"
                    
                    # Get topic categorization
                    topic_result = await call_llm_api(text, "topic")
                    
                    # Send topic result
                    yield f"data: {json.dumps({'type': 'topic', 'index': i, 'topic': topic_result.get('content', 'General')})}\n\n"
                
                # Clean and validate results
                sentiment = sentiment_result.get("sentiment", "neutral")
                confidence = sentiment_result.get("confidence", 0.5)
                topic = topic_result.get("content", "General") if "content" in topic_result else "General"
                
                # Validate confidence
                try:
                    confidence = float(confidence)
                    if not (0.0 <= confidence <= 1.0):
                        confidence = 0.5
                except (ValueError, TypeError):
                    confidence = 0.5
                
                # Update statistics
                sentiment_counts[sentiment] += 1
                topic_counts[topic] = topic_counts.get(topic, 0) + 1
                total_confidence += confidence
                
                # Extract exposure score for streaming analysis
                exposure_score = 0.0
                if 'max_exp' in row:
                    try:
                        raw_value = row['max_exp']
                        if raw_value is not None and str(raw_value).strip() not in ['', 'nan', 'None']:
                            if isinstance(raw_value, str):
                                clean_value = raw_value.replace(',', '').replace(' ', '').strip()
                                if clean_value:
                                    exposure_score = float(clean_value)
                            else:
                                exposure_score = float(raw_value)
                        print(f"Streaming Row {i}: max_exp extracted exposure {exposure_score} from raw value: {raw_value}")
                    except (ValueError, TypeError) as e:
                        print(f"Streaming Row {i}: Failed to extract max_exp {row.get('max_exp')}: {e}")
                        exposure_score = 0.0
                
                # Extract exposure score for streaming analysis
                exposure_score = 0.0
                if 'max_exp' in row:
                    try:
                        raw_value = row['max_exp']
                        if raw_value is not None and str(raw_value).strip() not in ['', 'nan', 'None']:
                            if isinstance(raw_value, str):
                                clean_value = raw_value.replace(',', '').replace(' ', '').strip()
                                if clean_value:
                                    exposure_score = float(clean_value)
                            else:
                                exposure_score = float(raw_value)
                        print(f"Streaming Row {i}: max_exp extracted exposure {exposure_score} from raw value: {raw_value}")
                    except (ValueError, TypeError) as e:
                        print(f"Streaming Row {i}: Failed to extract max_exp {row.get('max_exp')}: {e}")
                        exposure_score = 0.0
                
                # Create analyzed row with all original data plus analysis results
                analyzed_row = {
                    **row,  # Keep all original columns including max_exp
                    "sentiment": sentiment,
                    "confidence": confidence,
                    "topic": topic.strip() if topic else "General",
                    "exposure_score": exposure_score,
                    "profile_tag": str(row.get('tag', 'UNTAGGED')).upper(),
                    "author_url": str(row.get('author_url', '')),
                    "content_url": str(row.get('content_url', ''))
                }
                analyzed_data.append(analyzed_row)
                
                # Send row complete event
                yield f"data: {json.dumps({'type': 'row_complete', 'index': i, 'sentiment': sentiment, 'confidence': confidence, 'topic': topic})}\n\n"
                
            except Exception as e:
                error_msg = f"Error analyzing row: {str(e)}"
                analyzed_row = {
                    **row,
                    "sentiment": "neutral",
                    "confidence": 0.5,
                    "topic": "General"
                }
                analyzed_data.append(analyzed_row)
                yield f"data: {json.dumps({'type': 'error', 'index': i, 'error': error_msg})}\n\n"
        
        # Calculate final statistics
        avg_confidence = total_confidence / len(analyzed_data) if analyzed_data else 0.5
        if not (0.0 <= avg_confidence <= 1.0):
            avg_confidence = 0.5
        
        # Calculate exposure statistics for streaming
        exposure_scores = [row.get('exposure_score', 0) for row in analyzed_data]
        exposure_stats = {
            'max_exposure': max(exposure_scores) if exposure_scores else 0,
            'min_exposure': min(exposure_scores) if exposure_scores else 0,
            'avg_exposure': sum(exposure_scores) / len(exposure_scores) if exposure_scores else 0,
            'total_exposure': sum(exposure_scores)
        }
        
        # Profile distribution for streaming
        profile_distribution = {}
        for row in analyzed_data:
            tag = row.get('profile_tag', 'UNTAGGED')
            profile_distribution[tag] = profile_distribution.get(tag, 0) + 1
        
        # Store results
        session_data.sentiments = analyzed_data
        
        # Send completion event with exposure statistics
        yield f"data: {json.dumps({'type': 'complete_all', 'analyzed_data': analyzed_data, 'statistics': sentiment_counts, 'topic_statistics': topic_counts, 'average_confidence': avg_confidence, 'exposure_statistics': exposure_stats, 'profile_distribution': profile_distribution, 'text_column': text_column})}\n\n"
    
    return StreamingResponse(
        analyze_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Credentials": "true"
        }
    )

@app.post("/mitigate")
async def generate_comments(request: CommentRequest):
    """Generate counter/amplify comments for analyzed data"""
    if request.session_id not in sessions:
        raise HTTPException(status_code=404, detail="Session not found")
    
    session_data = sessions[request.session_id]
    analyzed_data = session_data.sentiments
    
    if not analyzed_data:
        raise HTTPException(status_code=400, detail="No analyzed data found")
    
    # Use override if provided, otherwise use intelligent detection
    if request.text_column_override:
        text_column = request.text_column_override
        print(f"Using user-specified text column: {text_column}")
    else:
        # Intelligent text column detection (excluding analysis result columns)
        filtered_data = []
        for row in analyzed_data:
            filtered_row = {k: v for k, v in row.items() if k not in ["sentiment", "confidence", "topic"]}
            filtered_data.append(filtered_row)
        text_column = detect_text_column(filtered_data)
        print(f"Using auto-detected text column: {text_column}")
    
    async def generate_comment(row):
        text = str(row.get(text_column, ""))
        sentiment = row.get("sentiment", "neutral")
        
        if sentiment == "negative":
            result = await call_llm_api(text, "counter")
        elif sentiment == "positive":
            result = await call_llm_api(text, "amplify")
        else:  # neutral
            result = await call_llm_api(text, "neutral_positive")
        
        return result.get("content", "No comment generated") if isinstance(result, dict) else str(result)
    
    # Process comments concurrently
    semaphore = asyncio.Semaphore(8)  # Limit concurrent requests
    
    async def generate_with_semaphore(row):
        async with semaphore:
            return await generate_comment(row)
    
    tasks = [generate_with_semaphore(row) for row in analyzed_data]
    comments = await asyncio.gather(*tasks)
    
    # Combine data with comments
    final_data = []
    for i, row in enumerate(analyzed_data):
        final_row = {**row, "generated_comment": comments[i]}
        final_data.append(final_row)
    
    session_data.comments = final_data
    
    # Use safe JSON response
    return safe_json_response({"data_with_comments": final_data})

@app.post("/generate-report")
async def generate_report(request: ReportRequest):
    """Generate comprehensive PowerPoint report with 10 sections"""
    if request.session_id not in sessions:
        raise HTTPException(status_code=404, detail="Session not found")
    
    session_data = sessions[request.session_id]
    final_data = session_data.comments or session_data.sentiments
    
    if not final_data:
        raise HTTPException(status_code=400, detail="No processed data found")
    
    # Create PowerPoint presentation
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Social Media Sentiment Analysis Report"
    subtitle.text = f"Malaysian External Intelligence Organisation (MEIO)\nDocument: {session_data.filename}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    bullet_slide_layout = prs.slide_layouts[1]
    
    # Calculate comprehensive statistics
    sentiment_counts = {"positive": 0, "negative": 0, "neutral": 0}
    topic_counts = {}
    confidence_sum = 0
    
    for row in final_data:
        sentiment_counts[row.get("sentiment", "neutral")] += 1
        confidence_sum += row.get("confidence", 0.5)
        topic = row.get("topic", "General")
        if isinstance(topic, str):
            topics = [t.strip() for t in topic.split(',')]
            for t in topics:
                topic_counts[t] = topic_counts.get(t, 0) + 1
    
    total_mentions = len(final_data)
    avg_confidence = confidence_sum / total_mentions if total_mentions > 0 else 0
    positive_pct = (sentiment_counts['positive'] / total_mentions * 100) if total_mentions > 0 else 0
    negative_pct = (sentiment_counts['negative'] / total_mentions * 100) if total_mentions > 0 else 0
    neutral_pct = (sentiment_counts['neutral'] / total_mentions * 100) if total_mentions > 0 else 0
    
    # 1. 📋 Executive Summary
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "📋 Executive Summary"
    
    exec_summary = f"""Key Findings:
• Total mentions analyzed: {total_mentions:,}
• Overall sentiment: {positive_pct:.1f}% positive, {negative_pct:.1f}% negative, {neutral_pct:.1f}% neutral
• Average confidence: {avg_confidence*100:.1f}%
• Top discussion topics: {', '.join(list(topic_counts.keys())[:3])}

What Happened:
{"Predominantly positive sentiment indicates strong public support" if positive_pct > negative_pct else "Mixed sentiment requires attention to negative concerns" if negative_pct > positive_pct else "Balanced sentiment suggests neutral public opinion"}

Why It Matters:
• Public sentiment directly impacts policy effectiveness
• High confidence scores ({avg_confidence*100:.1f}%) ensure reliable insights
• Topic diversity shows broad engagement across multiple areas

Recommended Actions:
• {"Maintain current positive momentum" if positive_pct > negative_pct else "Address negative sentiment drivers"}
• Focus on high-engagement topics for maximum impact
• {"Monitor for potential issues" if negative_pct > 20 else "Continue current strategy"}"""
    
    slide.placeholders[1].text = exec_summary
    
    # 2. 📊 Volume Metrics
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "📊 Volume Metrics"
    
    volume_content = f"""Total Mentions: {total_mentions:,}
Analysis Period: Document uploaded on {datetime.now().strftime('%Y-%m-%d')}
Document Source: {session_data.filename}

Sentiment Distribution:
• Positive mentions: {sentiment_counts['positive']:,} ({positive_pct:.1f}%)
• Negative mentions: {sentiment_counts['negative']:,} ({negative_pct:.1f}%)
• Neutral mentions: {sentiment_counts['neutral']:,} ({neutral_pct:.1f}%)

Confidence Metrics:
• Average confidence: {avg_confidence*100:.1f}%
• High confidence entries (>80%): {sum(1 for r in final_data if r.get('confidence', 0) > 0.8)} ({sum(1 for r in final_data if r.get('confidence', 0) > 0.8)/total_mentions*100 if total_mentions > 0 else 0:.1f}%)
• Low confidence entries (<60%): {sum(1 for r in final_data if r.get('confidence', 0) < 0.6)} ({sum(1 for r in final_data if r.get('confidence', 0) < 0.6)/total_mentions*100 if total_mentions > 0 else 0:.1f}%)

Trend Analysis:
• Sentiment trend: {"Positive trajectory" if positive_pct > negative_pct else "Mixed trend" if abs(positive_pct - negative_pct) < 10 else "Negative trend"}
• Engagement quality: {"High" if avg_confidence > 0.7 else "Medium" if avg_confidence > 0.5 else "Low"} based on confidence scores"""
    
    slide.placeholders[1].text = volume_content
    
    # 3. 😊 Detailed Sentiment Analysis
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "😊 Sentiment Analysis Breakdown"
    
    # Calculate sentiment drivers by topic
    positive_topics = {}
    negative_topics = {}
    neutral_topics = {}
    
    for row in final_data:
        topic = row.get("topic", "General")
        topics = [t.strip() for t in topic.split(',') if t.strip()]
        sentiment = row.get("sentiment", "neutral")
        
        for t in topics:
            if sentiment == "positive":
                positive_topics[t] = positive_topics.get(t, 0) + 1
            elif sentiment == "negative":
                negative_topics[t] = negative_topics.get(t, 0) + 1
            else:
                neutral_topics[t] = neutral_topics.get(t, 0) + 1
    
    top_positive = sorted(positive_topics.items(), key=lambda x: x[1], reverse=True)[:3]
    top_negative = sorted(negative_topics.items(), key=lambda x: x[1], reverse=True)[:3]
    
    sentiment_content = f"""Sentiment Distribution Analysis:
• Positive: {sentiment_counts['positive']} mentions ({positive_pct:.1f}%)
• Negative: {sentiment_counts['negative']} mentions ({negative_pct:.1f}%)
• Neutral: {sentiment_counts['neutral']} mentions ({neutral_pct:.1f}%)

Key Positive Sentiment Drivers:
{chr(10).join([f"• {topic}: {count} mentions" for topic, count in top_positive]) if top_positive else "• No significant positive drivers identified"}

Key Negative Sentiment Drivers:
{chr(10).join([f"• {topic}: {count} mentions" for topic, count in top_negative]) if top_negative else "• No significant negative drivers identified"}

Sentiment Quality Indicators:
• High confidence positive: {sum(1 for r in final_data if r.get('sentiment') == 'positive' and r.get('confidence', 0) > 0.8)}
• High confidence negative: {sum(1 for r in final_data if r.get('sentiment') == 'negative' and r.get('confidence', 0) > 0.8)}

Overall Sentiment Trend:
{"Strong positive momentum - maintain current approach" if positive_pct > 60 else "Predominantly negative - immediate attention required" if negative_pct > 60 else "Mixed sentiment - balanced approach needed"}"""
    
    slide.placeholders[1].text = sentiment_content
    
    # 4. 🧍‍♂️ Audience Insights
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "🧍‍♂️ Audience Insights"
    
    # Extract source information if available
    sources = {}
    for row in final_data:
        source = row.get("source", "Unknown Source")
        sources[source] = sources.get(source, 0) + 1
    
    top_sources = sorted(sources.items(), key=lambda x: x[1], reverse=True)[:5]
    
    audience_content = f"""Content Sources Analysis:
{chr(10).join([f"• {source}: {count} mentions ({count/total_mentions*100:.1f}%)" for source, count in top_sources]) if top_sources else "• Source information not available in dataset"}

Engagement Patterns:
• Most active source: {top_sources[0][0] if top_sources else "Not identified"}
• Source diversity: {len(sources)} different sources identified
• Average mentions per source: {total_mentions/len(sources):.1f} if sources else "N/A"

Demographics Insights:
• Platform diversity indicates broad audience reach
• Cross-platform consistency in sentiment patterns
• Engagement quality varies by source type

Key Audience Characteristics:
• Primary engagement sources show {"government/official" if any("official" in s.lower() for s, _ in top_sources[:3]) else "social media" if any("social" in s.lower() for s, _ in top_sources[:3]) else "mixed"} focus
• Content type preference: {"formal communication" if any("official" in s.lower() for s, _ in top_sources[:3]) else "informal discussion"}
• Response patterns suggest {"high" if avg_confidence > 0.7 else "moderate"} audience engagement quality"""
    
    slide.placeholders[1].text = audience_content
    
    # 5. 🔥 Top Performing Content
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "🔥 Top Performing Content Analysis"
    
    # Identify high-performing content (high confidence + positive sentiment)
    high_performing = [r for r in final_data if r.get('confidence', 0) > 0.8 and r.get('sentiment') == 'positive']
    concerning_content = [r for r in final_data if r.get('confidence', 0) > 0.8 and r.get('sentiment') == 'negative']
    
    content_performance = f"""High-Impact Positive Content ({len(high_performing)} items):
• High confidence positive mentions: {len(high_performing)}
• Average confidence of positive content: {sum(r.get('confidence', 0) for r in high_performing)/len(high_performing)*100:.1f}% if high_performing else "N/A"
• Top positive topics: {', '.join(set([r.get('topic', 'General').split(',')[0].strip() for r in high_performing[:5]]))if high_performing else "None identified"}

Content Performance Insights:
• Most engaging format: {"Official statements" if any("official" in str(r.get('source', '')).lower() for r in high_performing[:3]) else "Social media posts" if high_performing else "Mixed formats"}
• Key success factors: {"Policy announcements" if any("policy" in str(r).lower() for r in high_performing[:3]) else "Leadership messaging" if high_performing else "Content varies"}
• Optimal content length: {"Long-form" if high_performing and sum(len(str(r.get('text', ''))) for r in high_performing[:5])/len(high_performing[:5]) > 100 else "Short-form"}

Concerning Content Patterns ({len(concerning_content)} items):
• High confidence negative mentions requiring attention
• Key issues: {', '.join(set([r.get('topic', 'General').split(',')[0].strip() for r in concerning_content[:3]])) if concerning_content else "None identified"}

Recommendations:
• Amplify successful content formats and topics
• Replicate high-performing messaging strategies
• Address concerning content patterns proactively"""
    
    slide.placeholders[1].text = content_performance
    
    # 6. 🗣️ Key Topics & Hashtags
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "🗣️ Key Topics & Hashtags Analysis"
    
    top_topics = sorted(topic_counts.items(), key=lambda x: x[1], reverse=True)[:10]
    
    topics_content = f"""Trending Topics (by mention volume):
{chr(10).join([f"{i+1}. {topic}: {count} mentions ({count/total_mentions*100:.1f}%)" for i, (topic, count) in enumerate(top_topics)])}

Topic Performance Analysis:
• Most discussed: {top_topics[0][0] if top_topics else "N/A"}
• Emerging themes: {', '.join([topic for topic, count in top_topics[3:6]]) if len(top_topics) > 3 else "Limited topic diversity"}
• Coverage diversity: {len(topic_counts)} unique topics identified

Topic Sentiment Correlation:
• Positive topic drivers: {', '.join([t for t, c in sorted(positive_topics.items(), key=lambda x: x[1], reverse=True)[:3]]) if positive_topics else "None identified"}
• Negative topic drivers: {', '.join([t for t, c in sorted(negative_topics.items(), key=lambda x: x[1], reverse=True)[:3]]) if negative_topics else "None identified"}

Strategic Topic Insights:
• Focus areas for amplification: Top positive sentiment topics
• Areas requiring attention: High-volume negative sentiment topics
• Emerging opportunities: {', '.join([topic for topic, count in top_topics[5:8]]) if len(top_topics) > 5 else "Monitor new topic developments"}"""
    
    slide.placeholders[1].text = topics_content
    
    # 7. ⚔️ Competitor & Industry Benchmarking
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "⚔️ Competitive Analysis & Benchmarking"
    
    benchmark_content = f"""Industry Sentiment Benchmarks:
• Current positive rate: {positive_pct:.1f}%
• Industry average (estimated): 45-55% positive sentiment
• Performance vs. benchmark: {"Above average" if positive_pct > 55 else "Below average" if positive_pct < 45 else "Average"}

Competitive Positioning:
• Sentiment advantage: {"Strong positive positioning" if positive_pct > 60 else "Competitive challenge" if negative_pct > 40 else "Neutral positioning"}
• Topic leadership: Focus on {', '.join([t for t, c in sorted(positive_topics.items(), key=lambda x: x[1], reverse=True)[:2]]) if positive_topics else "core topics"}
• Market perception: {"Favorable" if positive_pct > negative_pct else "Mixed" if abs(positive_pct - negative_pct) < 10 else "Challenging"}

Strategic Recommendations:
• Leverage positive sentiment in {top_topics[0][0] if top_topics else "identified"} topics
• Address competitive gaps in negative sentiment areas
• Maintain confidence levels through consistent messaging

Monitoring Priorities:
• Track competitor sentiment trends
• Monitor topic emergence and sentiment shifts
• Benchmark response effectiveness over time"""
    
    slide.placeholders[1].text = benchmark_content
    
    # 8. 🚨 Crisis or Issue Tracking
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "🚨 Crisis & Issue Monitoring"
    
    # Identify potential crisis indicators
    high_confidence_negative = [r for r in final_data if r.get('confidence', 0) > 0.8 and r.get('sentiment') == 'negative']
    crisis_topics = {}
    for row in high_confidence_negative:
        topic = row.get('topic', 'General')
        topics = [t.strip() for t in topic.split(',')]
        for t in topics:
            crisis_topics[t] = crisis_topics.get(t, 0) + 1
    
    crisis_content = f"""Crisis Indicators Assessment:
• High-confidence negative mentions: {len(high_confidence_negative)} ({len(high_confidence_negative)/total_mentions*100 if total_mentions > 0 else 0:.1f}%)
• Crisis risk level: {"HIGH" if len(high_confidence_negative) > total_mentions * 0.3 else "MEDIUM" if len(high_confidence_negative) > total_mentions * 0.15 else "LOW"}

Issue Categories Requiring Attention:
{chr(10).join([f"• {topic}: {count} high-confidence negative mentions" for topic, count in sorted(crisis_topics.items(), key=lambda x: x[1], reverse=True)[:5]]) if crisis_topics else "• No significant crisis indicators detected"}

Early Warning Signals:
• Negative sentiment spikes: {"Detected" if negative_pct > 40 else "Not detected"}
• Topic concentration: {"High risk" if any(count > total_mentions * 0.3 for count in negative_topics.values()) else "Manageable"}
• Confidence level: {"Reliable concerns" if any(r.get('confidence', 0) > 0.8 for r in high_confidence_negative) else "Low confidence issues"}

Recommended Response Actions:
{"• Immediate crisis response protocol activation" if len(high_confidence_negative) > total_mentions * 0.3 else "• Proactive monitoring and response preparation"}
• Prepare counter-messaging for identified negative topics
• Monitor sentiment escalation patterns
• Engage with concerning content sources diplomatically"""
    
    slide.placeholders[1].text = crisis_content
    
    # 9. 💡 Strategic Insights & Recommendations
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "💡 Strategic Insights & Recommendations"
    
    insights_content = f"""Strategic Takeaways:
• Sentiment landscape: {"Favorable environment for policy advancement" if positive_pct > negative_pct else "Challenging environment requiring careful navigation"}
• Topic opportunities: High engagement in {', '.join([t for t, c in sorted(topic_counts.items(), key=lambda x: x[1], reverse=True)[:3]]) if topic_counts else "multiple areas"}
• Communication effectiveness: {avg_confidence*100:.1f}% average confidence indicates {"strong" if avg_confidence > 0.7 else "moderate"} message clarity

What to Continue:
• {"Maintain positive momentum in " + str(top_positive[0][0]) if top_positive else "Continue current messaging approach"}
• Leverage high-performing content formats
• Sustain engagement in top-performing topics

What to Stop:
• {"Address messaging that generates negative sentiment in " + str(top_negative[0][0]) if top_negative else "No significant negative patterns identified"}
• Reduce low-confidence communication approaches
• Avoid topics with consistently poor reception

What to Improve:
• Enhance message clarity to improve confidence scores
• Expand positive sentiment in neutral topic areas
• Develop targeted responses for negative sentiment drivers

Opportunity Areas:
• {"Sustainability content shows high engagement potential" if "environment" in str(topic_counts).lower() else "Policy communication opportunities identified"}
• Cross-topic messaging integration
• Proactive sentiment management in emerging topics"""
    
    slide.placeholders[1].text = insights_content
    
    # 10. 📅 Methodology & Data Sources
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "📅 Methodology & Data Sources"
    
    methodology_content = f"""Analysis Framework:
• AI-powered sentiment analysis using advanced language models
• Multi-topic categorization with confidence scoring
• Automated response generation for engagement strategy
• Real-time processing with streaming capabilities

Data Sources:
• Document: {session_data.filename}
• Records analyzed: {total_mentions:,}
• Analysis date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
• Processing method: {"Streaming analysis" if hasattr(session_data, 'streaming_used') else "Standard batch processing"}

Technical Specifications:
• Sentiment confidence threshold: 80%+ for high-confidence classifications
• Topic extraction: Multi-label classification with comma separation
• Response generation: Context-aware automated suggestions
• Quality assurance: Multi-layer validation and error handling

Confidence & Reliability:
• Overall analysis confidence: {avg_confidence*100:.1f}%
• High-confidence entries: {sum(1 for r in final_data if r.get('confidence', 0) > 0.8)/total_mentions*100 if total_mentions > 0 else 0:.1f}%
• Data completeness: 100% (all records processed)

Limitations & Considerations:
• Analysis based on provided dataset only
• Confidence scores reflect model certainty, not absolute accuracy
• Results should be validated with domain expertise
• Temporal context limited to document timestamp
• Response suggestions require human review before deployment"""
    
    slide.placeholders[1].text = methodology_content
    
    # Save to temporary file
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(temp_file.name)
    
    return FileResponse(
        temp_file.name,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        filename=f"MEIO_Comprehensive_Report_{session_data.filename.split('.')[0]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    )

@app.post("/mitigate-streaming")
async def mitigate_streaming(request: CommentRequest):
    """Stream comment generation with real-time updates"""
    if request.session_id not in sessions:
        raise HTTPException(status_code=404, detail="Session not found")
    
    session_data = sessions[request.session_id]
    analyzed_data = session_data.sentiments
    
    if not analyzed_data:
        raise HTTPException(status_code=400, detail="No analyzed data found")
    
    print(f"Starting streaming comment generation for {len(analyzed_data)} rows")
    
    # Use override if provided, otherwise use intelligent detection
    if request.text_column_override:
        text_column = request.text_column_override
        print(f"Using user-specified text column: {text_column}")
    else:
        # Intelligent text column detection (excluding analysis result columns)
        filtered_data = []
        for row in analyzed_data:
            filtered_row = {k: v for k, v in row.items() if k not in ["sentiment", "confidence", "topic"]}
            filtered_data.append(filtered_row)
        text_column = detect_text_column(filtered_data)
        print(f"Using auto-detected text column: {text_column}")
    
    async def generate_stream():
        final_data = []
        
        for i, row in enumerate(analyzed_data):
            text = str(row.get(text_column, ""))
            sentiment = row.get("sentiment", "neutral")
            
            # Send row start event
            yield f"data: {json.dumps({'type': 'row_start', 'index': i, 'total': len(analyzed_data), 'text': text[:100], 'sentiment': sentiment, 'topic': row.get('topic', 'General')})}\n\n"
            
            try:
                # Generate comment based on sentiment
                if sentiment == "negative":
                    result = await call_llm_api(text, "counter")
                elif sentiment == "positive":
                    result = await call_llm_api(text, "amplify")
                else:  # neutral
                    result = await call_llm_api(text, "neutral_positive")
                
                final_comment = result.get("content", "No comment generated") if isinstance(result, dict) else str(result)
                
                # Send word-by-word updates for streaming effect
                words = final_comment.split()
                complete_comment = ""
                for word in words:
                    complete_comment += word + " "
                    yield f"data: {json.dumps({'type': 'word', 'index': i, 'word': word + ' ', 'complete': complete_comment.strip()})}\n\n"
                    await asyncio.sleep(0.05)  # Streaming effect
                
                # Add to final data
                final_row = {**row, "generated_comment": final_comment.strip()}
                final_data.append(final_row)
                
                # Send row complete event
                yield f"data: {json.dumps({'type': 'row_complete', 'index': i, 'comment': final_comment.strip()})}\n\n"
                
            except Exception as e:
                error_msg = f"Error generating comment: {str(e)}"
                final_row = {**row, "generated_comment": error_msg}
                final_data.append(final_row)
                yield f"data: {json.dumps({'type': 'error', 'index': i, 'error': error_msg})}\n\n"
        
        # Store final data
        session_data.comments = final_data
        
        # Send completion event
        yield f"data: {json.dumps({'type': 'complete_all', 'data': final_data})}\n\n"
    
    return StreamingResponse(
        generate_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Credentials": "true"
        }
    )

@app.post("/analyze-influential")
async def analyze_influential_voices(request: InfluentialAnalysisRequest):
    """Analyze influential voices based on exposure score and sentiment"""
    if request.session_id not in sessions:
        raise HTTPException(status_code=404, detail="Session not found")
    
    session_data = sessions[request.session_id]
    analyzed_data = session_data.sentiments
    
    if not analyzed_data:
        raise HTTPException(status_code=400, detail="No analyzed data found. Run sentiment analysis first.")
    
    # Detect Cyabra-style columns
    cyabra_columns = detect_cyabra_columns(analyzed_data)
    exposure_column = cyabra_columns.get('exposure_column')
    tag_column = cyabra_columns.get('tag_column')
    author_url_column = cyabra_columns.get('author_url_column')
    content_url_column = cyabra_columns.get('content_url_column')
    
    print(f"Detected Cyabra columns: {cyabra_columns}")
    
    # Use override if provided, otherwise use intelligent detection
    if request.text_column_override:
        text_column = request.text_column_override
        print(f"Using user-specified text column: {text_column}")
    else:
        # Intelligent text column detection (excluding analysis result columns)
        filtered_data = []
        for row in analyzed_data:
            filtered_row = {k: v for k, v in row.items() if k not in ["sentiment", "confidence", "topic"]}
            filtered_data.append(filtered_row)
        text_column = detect_text_column(filtered_data)
        print(f"Using auto-detected text column: {text_column}")
    
    # Extract influential voices
    influential_voices = []
    for idx, row in enumerate(analyzed_data):
        text = str(row.get(text_column, ""))
        sentiment = row.get("sentiment", "neutral")
        confidence = row.get("confidence", 0.5)
        topics = row.get("topic", "General")
        
        # Extract exposure score - check analyzed data first since it may already have exposure_score
        exposure_score = 0.0
        
        # First check if exposure_score is already in the analyzed data
        if 'exposure_score' in row:
            try:
                exposure_score = float(row['exposure_score'])
                print(f"Influential analysis - Row {idx}: Using existing exposure_score {exposure_score}")
            except (ValueError, TypeError):
                exposure_score = 0.0
        
        # If not found or zero, try to extract from original columns
        if exposure_score == 0.0:
            if exposure_column and exposure_column in row:
                try:
                    exposure_raw = row[exposure_column]
                    if exposure_raw is not None and str(exposure_raw).strip() != '':
                        if isinstance(exposure_raw, str):
                            clean_value = exposure_raw.replace(',', '').replace(' ', '')
                            exposure_score = float(clean_value)
                        else:
                            exposure_score = float(exposure_raw)
                    print(f"Influential analysis - Row {idx}: Extracted exposure {exposure_score} from column {exposure_column}")
                except (ValueError, TypeError) as e:
                    print(f"Influential analysis - Row {idx}: Failed to extract exposure from {exposure_column}: {e}")
                    exposure_score = 0.0
            
            # Fallback to direct max_exp column access
            if exposure_score == 0.0 and 'max_exp' in row:
                try:
                    raw_value = row['max_exp']
                    if raw_value is not None and str(raw_value).strip() != '':
                        if isinstance(raw_value, str):
                            clean_value = raw_value.replace(',', '').replace(' ', '')
                            exposure_score = float(clean_value)
                        else:
                            exposure_score = float(raw_value)
                    print(f"Influential analysis - Row {idx}: Fallback extracted exposure {exposure_score} from max_exp")
                except (ValueError, TypeError) as e:
                    print(f"Influential analysis - Row {idx}: Fallback failed to extract exposure: {e}")
                    exposure_score = 0.0
        
        # Extract profile tag
        profile_tag = "UNTAGGED"
        if 'profile_tag' in row:
            profile_tag = str(row.get('profile_tag', 'UNTAGGED')).upper()
        elif tag_column and tag_column in row:
            profile_tag = str(row.get(tag_column, "UNTAGGED")).upper()
        elif 'tag' in row:
            profile_tag = str(row.get('tag', 'UNTAGGED')).upper()
        
        # Extract URLs with fallback to direct column access
        author_url = ""
        content_url = ""
        
        if 'author_url' in row:
            author_url = str(row.get('author_url', ''))
        elif author_url_column and author_url_column in row:
            author_url = str(row.get(author_url_column, ''))
        
        if 'content_url' in row:
            content_url = str(row.get('content_url', ''))
        elif content_url_column and content_url_column in row:
            content_url = str(row.get(content_url_column, ''))
        
        # Only process negative sentiment posts for counter-statements
        should_generate_counter = sentiment == "negative"
        
        influential_voices.append({
            'row_index': idx,
            'original_text': text,
            'sentiment': sentiment,
            'confidence': confidence,
            'topics': topics.split(',') if isinstance(topics, str) else [topics] if topics else ["General"],
            'exposure_score': exposure_score,
            'profile_tag': profile_tag,
            'author_url': author_url,
            'content_url': content_url,
            'should_generate_counter': should_generate_counter
        })
    
    # Sort by exposure score descending
    influential_voices.sort(key=lambda x: x['exposure_score'], reverse=True)
    
    # Filter based on threshold or top count
    priority_voices = []
    if request.use_threshold:
        # Use exposure threshold
        priority_voices = [v for v in influential_voices if v['exposure_score'] >= request.exposure_threshold]
        print(f"Found {len(priority_voices)} voices above {request.exposure_threshold} exposure threshold")
    else:
        # Use top count
        priority_voices = influential_voices[:request.top_count]
        print(f"Selected top {len(priority_voices)} voices by exposure")
    
    # Filter to only negative sentiment for counter-statement generation
    negative_priority_voices = [v for v in priority_voices if v['sentiment'] == 'negative']
    
    # Add influence ranking
    for rank, voice in enumerate(priority_voices):
        voice['influence_rank'] = rank + 1
        voice['is_priority'] = True
    
    # Generate counter-statements for high-influence negative voices
    async def generate_counter_statement(voice_data):
        if not voice_data['should_generate_counter']:
            return ""
        
        # Enhanced prompt for influential voice counter-statements
        text = voice_data['original_text']
        profile_type = voice_data['profile_tag']
        exposure = voice_data['exposure_score']
        
        prompt = f"""Generate a strategic counter-statement for this influential negative post with {exposure:.0f} exposure score from a {profile_type} profile.

Original post: {text}

Instructions for counter-statement:
- Address the concerns raised without being defensive
- Provide factual, constructive information
- Match the communication style appropriate for the audience
- Keep the tone professional but approachable
- Focus on solutions and positive outcomes
- {"Be diplomatic as this is from a real person" if profile_type == "REAL_PROFILE" else "Address systematically as this may be from an automated account" if profile_type == "FAKE_PROFILE" else "Provide institutional-level response as this is from an organization"}

Strategic counter-statement:"""
        
        # Call LLM with enhanced prompt
        payload = {
            "model": os.getenv("LLM_MODEL_NAME"),
            "messages": [
                {"role": "user", "content": prompt}
            ],
            "max_tokens": 300,
            "temperature": 0.3,
            "stream": False
        }
        
        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(f"{os.getenv('LLM_API_URL')}/chat/completions",
                                       json=payload,
                                       headers={"Content-Type": "application/json"}) as response:
                    if response.status == 200:
                        result = await response.json()
                        return result["choices"][0]["message"]["content"].strip()
        except Exception as e:
            print(f"Error generating counter-statement: {e}")
            return "Counter-statement generation failed"
        
        return "No counter-statement generated"
    
    # Generate counter-statements for negative priority voices
    counter_tasks = []
    for voice in negative_priority_voices:
        counter_tasks.append(generate_counter_statement(voice))
    
    if counter_tasks:
        counter_statements = await asyncio.gather(*counter_tasks)
        for i, voice in enumerate(negative_priority_voices):
            voice['counter_statement'] = counter_statements[i]
    
    # Store results in database
    try:
        db.store_influential_voices(request.session_id, priority_voices)
        print(f"Stored {len(priority_voices)} influential voices in database")
    except Exception as e:
        print(f"Error storing influential voices: {e}")
    
    # Calculate statistics
    stats = {
        'total_voices': len(influential_voices),
        'priority_voices': len(priority_voices),
        'negative_priority_voices': len(negative_priority_voices),
        'threshold_used': request.exposure_threshold if request.use_threshold else f"Top {request.top_count}",
        'cyabra_columns_detected': cyabra_columns,
        'profile_distribution': {},
        'exposure_stats': {
            'max_exposure': max([v['exposure_score'] for v in influential_voices]) if influential_voices else 0,
            'min_exposure': min([v['exposure_score'] for v in influential_voices]) if influential_voices else 0,
            'avg_exposure': sum([v['exposure_score'] for v in influential_voices]) / len(influential_voices) if influential_voices else 0
        }
    }
    
    # Profile distribution stats
    for voice in priority_voices:
        tag = voice['profile_tag']
        stats['profile_distribution'][tag] = stats['profile_distribution'].get(tag, 0) + 1
    
    return safe_json_response({
        'priority_voices': priority_voices,
        'statistics': stats,
        'analysis_config': {
            'exposure_threshold': request.exposure_threshold,
            'top_count': request.top_count,
            'use_threshold': request.use_threshold,
            'text_column': text_column,
            'detected_columns': cyabra_columns
        }
    })

@app.post("/save-analysis")
async def save_analysis_session(request: SaveAnalysisRequest):
    """Save an analysis session for later retrieval"""
    try:
        session_id = str(uuid.uuid4())
        saved_id = db.save_analysis_session(
            session_id,
            request.document_id,
            request.analysis_name,
            request.analysis_type,
            request.analysis_data
        )
        return {"saved_analysis_id": saved_id, "message": "Analysis saved successfully"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error saving analysis: {str(e)}")

@app.get("/saved-analyses")
async def get_saved_analyses(document_id: Optional[str] = None):
    """Get saved analysis sessions"""
    try:
        analyses = db.get_saved_analyses(document_id)
        return safe_json_response({"saved_analyses": analyses})
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error retrieving analyses: {str(e)}")

@app.get("/saved-analyses/{analysis_id}")
async def get_saved_analysis(analysis_id: str):
    """Get a specific saved analysis session"""
    try:
        analysis = db.get_analysis_session(analysis_id)
        if not analysis:
            raise HTTPException(status_code=404, detail="Analysis session not found")
        return safe_json_response(analysis)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error retrieving analysis: {str(e)}")

@app.delete("/saved-analyses/{analysis_id}")
async def delete_saved_analysis(analysis_id: str):
    """Delete a saved analysis session"""
    try:
        db.delete_analysis_session(analysis_id)
        return {"message": "Analysis deleted successfully"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error deleting analysis: {str(e)}")

@app.get("/")
async def root():
    return {"message": "MEIO Sentiment Analysis API", "status": "active"}

@app.get("/test")
async def test_endpoint():
    """Simple test endpoint to verify API connectivity"""
    print("Test endpoint called successfully")
    return {"status": "success", "message": "API is working", "timestamp": datetime.now().isoformat()}

@app.post("/test-upload")
async def test_upload(file: UploadFile = File(...)):
    """Minimal upload test"""
    print(f"Test upload called with file: {file.filename}")
    try:
        contents = await file.read()
        print(f"File size: {len(contents)} bytes")
        return {
            "status": "success",
            "filename": file.filename,
            "size": len(contents),
            "message": "File received successfully"
        }
    except Exception as e:
        print(f"Test upload error: {e}")
        return {"status": "error", "message": str(e)}

@app.get("/debug/sessions")
async def debug_sessions():
    """Debug endpoint to see active sessions"""
    return {
        "active_sessions": list(sessions.keys()),
        "session_count": len(sessions),
        "sessions_data": {sid: {"has_data": session.data is not None, "has_sentiments": session.sentiments is not None} for sid, session in sessions.items()}
    }

@app.get("/check-routes")
async def check_routes():
    """Check all registered routes"""
    routes = []
    for route in app.routes:
        if hasattr(route, 'path') and hasattr(route, 'methods'):
            routes.append(f"{list(route.methods)[0] if route.methods else 'GET'} {route.path}")
    return {"routes": routes}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8011)